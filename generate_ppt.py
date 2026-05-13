"""
生成论文智能助手项目演示 PPT
运行: python generate_ppt.py
输出: 论文智能助手-项目演示.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色方案 ──
NAVY   = RGBColor(0x0B, 0x1D, 0x3A)   # 深蓝背景
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xD4, 0xA8, 0x3C)   # 金色强调
LIGHT  = RGBColor(0xE8, 0xEC, 0xF0)   # 浅灰文字
BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # 亮蓝
GREEN  = RGBColor(0x10, 0xB9, 0x81)   # 成功绿
RED    = RGBColor(0xEF, 0x44, 0x44)   # 警告红
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)   # 橙色
CARD_BG = RGBColor(0x13, 0x2A, 0x4F)  # 卡片背景

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── 辅助函数 ──
def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT):
    """segments: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, seg in enumerate(segments):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = alignment
        text, size, color, bold = seg
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "微软雅黑"
    return tf

def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, width=0.06, height=0.8, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow_right(slide, left, top, width, height, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_divider(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_step_number(slide, left, top, number):
    circle = add_circle(slide, left, top, 0.45, GOLD)
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    return circle


# ═══════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# 装饰线
add_divider(slide, 1.5, 2.8, 2.0, GOLD)

add_text_box(slide, 1.5, 0.8, 10, 1.2,
    "论文智能助手", 56, WHITE, True, PP_ALIGN.LEFT)

add_text_box(slide, 1.5, 1.9, 10, 0.8,
    "基于多 Agent 协作与 RAG 检索增强的学术写作平台", 24, LIGHT, False, PP_ALIGN.LEFT)

add_text_box(slide, 1.5, 3.2, 10, 1.5,
    "从文献管理 → 数据分析 → 分章扩写 → 引用验证 → 一键导出\n不是简单问 ChatGPT，而是三个 AI 互相审查的完整学术工作流",
    18, LIGHT, False, PP_ALIGN.LEFT)

add_divider(slide, 1.5, 5.0, 2.0, GOLD)
add_text_box(slide, 1.5, 5.3, 10, 0.6,
    "农业资源与环境实验室  ·  2026", 16, RGBColor(0x88, 0x99, 0xAA), False)


# ═══════════════════════════════════════
# Slide 2: 为什么做这个？— 痛点
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "写一篇 SCI 论文要经历什么？", 32, WHITE, True)

pain_points = [
    ("📚 查文献", "逐篇下载、分类、阅读\n动辄上百篇，耗时数周"),
    ("🔬 做实验 + 分析", "数据整理、图表制作\n趋势分析靠肉眼"),
    ("✍️ 写初稿", "组织语言、插入引用\n格式排版繁琐"),
    ("🔍 修改审查", "查重、降重、引文核对\n一致性检查容易遗漏"),
    ("📤 投稿排版", "不同期刊不同模板\n反复调整格式"),
]

for i, (title, desc) in enumerate(pain_points):
    x = 0.6 + i * 2.5
    add_card(slide, x, 1.6, 2.2, 2.6)
    add_text_box(slide, x + 0.2, 1.8, 1.8, 0.5, title, 18, WHITE, True)
    add_text_box(slide, x + 0.2, 2.4, 1.8, 1.5, desc, 13, LIGHT)

    if i < 4:
        add_text_box(slide, x + 2.3, 2.5, 0.3, 0.5, "→", 24, GOLD, True)

add_text_box(slide, 0.8, 4.8, 11.5, 1.5,
    "💡 我的目标：AI 帮你把机械劳动自动化，让你把时间花在真正需要思考的地方",
    20, GOLD, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 3: 功能总览 — 六步工作流
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "六大核心功能 — 覆盖论文写作全流程", 32, WHITE, True)

steps = [
    ("1", "建知识库", "批量上传PDF\n自动分类索引\nAI 文献对话"),
    ("2", "项目规划", "输入研究方向\nAI 生成大纲\n手动调整"),
    ("3", "数据分析+作图", "上传实验数据\nAI 趋势分析\nXRD/图表/流程图"),
    ("4", "分章扩写", "五章独立扩写\n三级子节编号\n自动配图插入"),
    ("5", "质量保障", "跨章一致性检查\n引用逐条验证\n查重+AI降重"),
    ("6", "一键导出", "Word/Markdown/PDF\nSCI/Nature/国标\n参考文献自动排版"),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.35 + i * 2.15
    add_card(slide, x, 1.5, 2.0, 3.5)
    add_text_box(slide, x + 0.15, 1.65, 0.5, 0.5, num, 28, GOLD, True)
    add_text_box(slide, x + 0.15, 2.15, 1.7, 0.5, title, 16, WHITE, True)
    add_divider(slide, x + 0.15, 2.65, 0.8, GOLD)
    add_text_box(slide, x + 0.15, 2.85, 1.7, 2.0, desc, 12, LIGHT)

    if i < 5:
        add_text_box(slide, x + 2.05, 2.8, 0.15, 0.5, "▸", 18, GOLD, True)

add_text_box(slide, 0.8, 5.5, 11.5, 1.2,
    "全部集成在一个工作台内，数据在功能间流转，无需切换工具",
    16, LIGHT, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 4: 核心特色 — 不是简单"问 ChatGPT"
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "核心特色：不是简单「问 ChatGPT」", 32, WHITE, True)

# 对比表格
add_card(slide, 0.5, 1.5, 5.8, 5.2, CARD_BG)
add_text_box(slide, 0.8, 1.7, 5.2, 0.5, "❌ 普通 ChatGPT 写论文", 20, RED, True)
compares_bad = [
    "靠训练数据「记忆」编造内容",
    "引用文献是虚构的，查不到原文",
    "一个 prompt 生成全部，无质量审查",
    "不懂你的研究方向，泛泛而谈",
    "图表？只能描述，画不出来",
    "不会帮你管理文献库",
]
for i, item in enumerate(compares_bad):
    add_text_box(slide, 1.0, 2.4 + i * 0.55, 5.0, 0.5, f"  · {item}", 14, LIGHT)

add_card(slide, 6.8, 1.5, 5.8, 5.2, RGBColor(0x0A, 0x3D, 0x1A))
add_text_box(slide, 7.1, 1.7, 5.2, 0.5, "✅ 本系统", 20, GREEN, True)
compares_good = [
    "基于你的 167 篇文献库检索后写作",
    "引用编号可追溯到文献原文段落",
    "三个 AI 模型互相审查（写→审→修）",
    "按研究方向精准检索相关文献",
    "8 种专业图表一键生成并插入论文",
    "统一文献管理 + 写作 + 导出工作台",
]
for i, item in enumerate(compares_good):
    add_text_box(slide, 7.3, 2.4 + i * 0.55, 5.0, 0.5, f"  · {item}", 14, LIGHT)


# ═══════════════════════════════════════
# Slide 5: 多 Agent 架构
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "科学依据 ① — 多 Agent 协作架构", 32, WHITE, True)
add_text_box(slide, 1.1, 1.15, 10, 0.5,
    "类比学术出版流程：自己写初稿 → 导师审稿提意见 → 修改后投稿", 16, LIGHT)

# 三个 Agent 卡片
agents = [
    ("Writer\n写手", "DeepSeek", "基于文献库\n起草初稿",
     "你自己写初稿", BLUE, "① 起草"),
    ("Verifier\n审稿人", "智谱 AI", "独立核查引用真实性\n逐条比对原文",
     "导师/审稿人", ORANGE, "② 审查"),
    ("Refiner\n主编", "条件触发", "根据审稿意见逐条修正\n严禁删引用逃避审查",
     "修改后定稿", GREEN, "③ 修正"),
]

for i, (name, model, duty, analogy, color, label) in enumerate(agents):
    x = 0.8 + i * 4.2
    add_card(slide, x, 2.0, 3.8, 3.8)
    # number badge
    add_circle(slide, x + 0.2, 2.2, 0.6, color)
    add_text_box(slide, x + 0.3, 2.25, 0.4, 0.5, label, 12, WHITE, True,
                 PP_ALIGN.CENTER)
    add_text_box(slide, x + 1.0, 2.2, 2.5, 0.6, name, 24, WHITE, True)
    add_text_box(slide, x + 0.3, 2.9, 1.5, 0.4, f"模型: {model}", 11, GOLD, False)
    add_divider(slide, x + 0.3, 3.4, 1.5, GOLD)
    add_text_box(slide, x + 0.3, 3.6, 3.2, 1.5, duty, 14, LIGHT)
    add_text_box(slide, x + 0.3, 4.7, 3.2, 0.5, f"≈ {analogy}", 13, RGBColor(0x88, 0x99, 0xAA), False)

# 论文引用
add_text_box(slide, 0.8, 6.2, 11.5, 0.8,
    "📄 学术依据: Chain-of-Verification (Dhuliawala et al., 2024) · Multi-Agent Debate (Du et al., 2023)",
    13, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 6: RAG 检索增强生成
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "科学依据 ② — RAG 检索增强生成", 32, WHITE, True)
add_text_box(slide, 1.1, 1.15, 10, 0.5,
    "不是让 AI 凭记忆编内容，而是先检索你的文献库，再基于真实文献写作", 16, LIGHT)

# 流程图
flow_items = [
    ("用户输入\n题目+方向+上下文", BLUE),
    ("混合检索\nBM25关键词 + 向量语义\n→ RRF融合排序", ORANGE),
    ("文献过滤\n按研究方向匹配分类\n来源去重(≤4条/篇)", GREEN),
    ("AI 写作\n基于 Top-20\n相关文献片段", BLUE),
]

for i, (text, color) in enumerate(flow_items):
    x = 0.5 + i * 3.2
    add_card(slide, x, 2.0, 2.8, 2.5, color if i % 2 == 0 else ORANGE if i == 1 else GREEN)
    add_text_box(slide, x + 0.2, 2.3, 2.4, 2.0, text, 15, WHITE, True, PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + 2.85, 2.9, 0.4, 0.5, "▸", 24, GOLD, True)

# 底部说明
add_card(slide, 0.5, 5.0, 12.3, 1.8, CARD_BG)
add_rich_text(slide, 0.8, 5.15, 11.8, 1.5, [
    ("混合检索策略", 16, GOLD, True),
    ("BM25 关键词检索：精确匹配实验术语和关键概念", 13, LIGHT, False),
    ("向量语义检索：理解查询意图，找到用词不同但意思相近的文献", 13, LIGHT, False),
    ("RRF (Reciprocal Rank Fusion)：融合两种排序信号，取各自优势  |  Cormack et al., 2009", 13, RGBColor(0x88,0x99,0xAA), False),
    ("章节感知 + 方向过滤：检索时自动注入章节类型关键词 + 研究方向分类，精准定位相关文献", 13, LIGHT, False),
])


# ═══════════════════════════════════════
# Slide 7: 引用真实性验证
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "科学依据 ③ — 引用真实性逐条验证", 32, WHITE, True)
add_text_box(slide, 1.1, 1.15, 10, 0.5,
    "ChatGPT 最严重的学术作弊：编造看起来很专业的 [1] [2]，但论文根本不存在", 16, LIGHT)

# 验证流程
verify_steps = [
    ("Writer 写初稿", "引用编号 [n] 必须对应\n文献库中真实存在的文献", BLUE),
    ("Verifier 拿原文比对", "获取被引用文献的完整原文\n逐句核实：观点/数据/结论", ORANGE),
    ("三种判定结果", "✓ 通过  ⚠ 归属错误  ✗ 疑似虚构", GREEN),
    ("Refiner 逐条修正", "严禁为通过审查而删除引用\n必须找到正确证据替换", BLUE),
]

for i, (title, desc, color) in enumerate(verify_steps):
    x = 0.5 + i * 3.2
    add_card(slide, x, 2.0, 2.8, 2.5)
    add_text_box(slide, x + 0.2, 2.15, 2.4, 0.5, title, 16, color, True)
    add_divider(slide, x + 0.2, 2.7, 1.5, color)
    add_text_box(slide, x + 0.2, 2.9, 2.4, 1.5, desc, 13, LIGHT)
    if i < 3:
        add_text_box(slide, x + 2.85, 2.9, 0.4, 0.5, "▸", 24, GOLD, True)

add_card(slide, 0.5, 5.0, 12.3, 1.8, CARD_BG)
add_rich_text(slide, 0.8, 5.15, 11.8, 1.5, [
    ("三层防护", 16, GOLD, True),
    ("① Writer 只能引用文献库中标记为「参考来源 [n]」的真实文献", 13, LIGHT, False),
    ("② Verifier 拿到被引用文献完整原文做事实对照，不是只看摘要", 13, LIGHT, False),
    ("③ 关键词重叠率辅助筛查：低于 8% → 标记疑似虚构 → 提醒人工核实", 13, LIGHT, False),
    ("NLI（自然语言推理）范式在学术写作中的应用  |  Bowman et al., 2015", 13, RGBColor(0x88,0x99,0xAA), False),
])


# ═══════════════════════════════════════
# Slide 8: Prompt 认知负载优化
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "科学依据 ④ — Prompt 认知负载优化", 32, WHITE, True)

add_card(slide, 0.5, 1.5, 5.8, 5.5, CARD_BG)
add_text_box(slide, 0.8, 1.7, 5.2, 0.5, "❌ 常见做法：一个 prompt 塞 9 条要求", 20, RED, True)
add_text_box(slide, 0.8, 2.3, 5.2, 4.5,
    """1. 语言规范
2. 深度结合文献库
3. 引用规范（6条子规则）
4. 术语准确
5. 逻辑连贯
6. 禁止标题
7. 禁止解释性文字
8. 段落格式要求
9. 配图要求

→ AI 注意力被稀释
   只认真看前 3-4 条
   中间部分基本忽略""", 14, LIGHT)

add_card(slide, 6.8, 1.5, 5.8, 5.5, RGBColor(0x0A, 0x3D, 0x1A))
add_text_box(slide, 7.1, 1.7, 5.2, 0.5, "✅ 本系统：3 组核心原则", 20, GREEN, True)
add_text_box(slide, 7.1, 2.3, 5.2, 4.5,
    """—— 核心写作原则 ——
原则1·学术质量（首因效应）
原则2·深度结合文献
原则3·结构与配图

—— 一致性约束 ——
（中间位置，数据参考）

—— ⚠️ 引用铁律（近因效应）——
· 严禁虚构引用
· 禁止输出元文字

→ 首因+近因效应
   AI 对关键规则的遵循度
   显著提升""", 14, LIGHT)

add_text_box(slide, 0.8, 7.0, 11.5, 0.5,
    "📄 Lost in the Middle (Liu et al., 2023) — 大模型对 prompt 中间部分的注意力衰减严重",
    12, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 9: 集成科学作图
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "集成科学作图 — 8 种专业图表", 32, WHITE, True)
add_text_box(slide, 1.1, 1.15, 10, 0.5,
    "不只是写文字 — 论文所需图表从生成到插入，全在平台内完成", 16, LIGHT)

charts = [
    ("📊", "通用图表", "折线/柱状/散点\n实验数据可视化", BLUE),
    ("🔬", "XRD 峰拟合", "自动分解重叠峰\n背景扣除+峰位标注", ORANGE),
    ("💎", "XRD 晶胞", "3D 晶体结构\n参数可视化", GREEN),
    ("📐", "Bragg 计算", "d值/2θ值\n晶面指数计算", BLUE),
    ("🔮", "非晶分析", "非晶相定量\n结晶度计算", ORANGE),
    ("🔄", "流程图", "文字描述→自动绘制\n实验流程图", GREEN),
    ("⚛️", "分子结构", "SMILES/分子式\n→化学结构式", BLUE),
    ("⚡", "机理图", "反应路径描述\n→自动绘制机理图", ORANGE),
]

for i, (icon, name, desc, color) in enumerate(charts):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.2
    y = 1.9 + row * 2.7
    add_card(slide, x, y, 2.9, 2.4)
    add_text_box(slide, x + 0.2, y + 0.15, 0.5, 0.5, icon, 24, WHITE, False)
    add_text_box(slide, x + 0.8, y + 0.15, 1.8, 0.5, name, 16, color, True)
    add_divider(slide, x + 0.2, y + 0.75, 1.5, color)
    add_text_box(slide, x + 0.2, y + 0.95, 2.5, 1.2, desc, 13, LIGHT)

add_text_box(slide, 0.8, 7.0, 11.5, 0.4,
    "产出高清 Base64/PNG 图片，一键插入论文正文，可直接用于投稿",
    13, LIGHT, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 10: 完整技术架构
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "完整技术架构 — 一页看懂全链路", 32, WHITE, True)

# 架构层级
layers = [
    ("用户层", "题目 + 章节 + 研究上下文 + 研究方向", BLUE),
    ("检索层", "BM25 关键词 + 向量语义 → RRF 融合 → 方向分类过滤 → 来源去重", ORANGE),
    ("生成层", "Writer (DeepSeek) → Verifier (智谱 AI) → Refiner (条件触发)", GREEN),
    ("后处理", "引用重排 + 层级编号 + 图表生成 + 一致性检查", BLUE),
    ("输出层", "Markdown 编辑 → SCI/Nature/国标 模板 → Word/PDF 导出", ORANGE),
]

for i, (name, desc, color) in enumerate(layers):
    y = 1.7 + i * 1.05
    add_card(slide, 1.0, y, 11.3, 0.9, color)
    add_text_box(slide, 1.2, y + 0.1, 2.0, 0.7, name, 18, WHITE, True)
    add_text_box(slide, 3.3, y + 0.15, 8.5, 0.6, desc, 15, WHITE, False)


# ═══════════════════════════════════════
# Slide 11: 工程亮点
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "工程亮点", 32, WHITE, True)

highlights = [
    ("167篇文献 12,626知识块", "本地 RAG 索引，无需联网检索", BLUE),
    ("混血检索 + 优雅降级", "Embedding API 不可用时，自动退化为 BM25 纯关键词检索", ORANGE),
    ("章节感知 + 方向过滤", "不同章节自动调整检索策略，研究方向精准分类", GREEN),
    ("引用逐条溯源", "每个 [n] 可追溯到文献原文段落", BLUE),
    ("多 Agent 独立模型", "Writer 用 DeepSeek，Verifier 用智谱，真正独立审查", ORANGE),
    ("认知负载优化", "9 条规则 → 3 组原则，首因+近因效应", GREEN),
    ("流式 SSE + 实时反馈", "生成过程实时可见，配图异步生成不阻塞写作", BLUE),
    ("多模板导出", "SCI / Nature / IEEE / GB/T 7713，参考文献自动排版", ORANGE),
]

for i, (title, desc, color) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.5 + row * 1.4
    add_card(slide, x, y, 6.0, 1.15)
    add_text_box(slide, x + 0.2, y + 0.1, 5.6, 0.45, title, 15, color, True)
    add_text_box(slide, x + 0.2, y + 0.55, 5.6, 0.45, desc, 13, LIGHT)


# ═══════════════════════════════════════
# Slide 12: 关键数据
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "关键数据", 32, WHITE, True)

stats = [
    ("167", "篇 PDF 文献", "已索引可检索"),
    ("12,626", "个知识块", "每块1000字"),
    ("3", "个 AI Agent", "写→审→修 协作"),
    ("8", "种专业图表", "XRD/流程图/分子图"),
    ("5", "种导出模板", "SCI/Nature/IEEE/国标"),
    ("0", "次引用虚构", "逐条原文验证"),
]

for i, (num, unit, sub) in enumerate(stats):
    x = 0.5 + i * 2.15
    add_card(slide, x, 2.0, 2.0, 3.5)
    add_text_box(slide, x + 0.15, 2.3, 1.7, 1.0, num, 48, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 3.4, 1.7, 0.5, unit, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 3.9, 1.7, 0.5, sub, 12, LIGHT, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 13: 未来规划
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "持续优化方向", 32, WHITE, True)

future_items = [
    ("近期", "🔴", [
        "RAG 检索支持 category 精准过滤（已实现）",
        "研究方向关键词注入检索 query（已实现）",
        "多方向混合文献库智能匹配",
    ]),
    ("中期", "🟡", [
        "向量存储从 JSON 迁移到 PGLite + HNSW 索引",
        "引入 NLI 模型做引用事实核查（替代关键词重叠）",
        "支持更多 LLM 后端（Claude, GPT-4o 等）",
    ]),
    ("远期", "🟢", [
        "迁移到专业向量数据库 (Qdrant/Milvus) 支撑百万级文献",
        "引用证据溯源：每个声明自动关联原文段落",
        "协同写作：多人同时编辑 + AI 辅助",
    ]),
]

for i, (period, emoji, items) in enumerate(future_items):
    x = 0.5 + i * 4.2
    add_card(slide, x, 1.5, 3.9, 5.2)
    add_text_box(slide, x + 0.2, 1.65, 3.5, 0.5, f"{emoji} {period}", 22, WHITE, True)
    add_divider(slide, x + 0.2, 2.2, 2.0, GOLD)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.2, 2.5 + j * 0.8, 3.5, 0.7, f"· {item}", 14, LIGHT)


# ═══════════════════════════════════════
# Slide 14: Q&A 预判
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_accent_bar(slide, 0.8, 0.6, height=0.5)
add_text_box(slide, 1.1, 0.55, 10, 0.7, "常见问题", 32, WHITE, True)

qas = [
    ("Q: 和 ChatGPT 有什么区别？",
     "A: ChatGPT 凭记忆编内容、引用虚构。本系统基于你的文献库检索后写作，三个 AI 互相审查，每个引用可追溯到原文段落。"),
    ("Q: AI 写的内容能直接投稿吗？",
     "A: 不能也不应该。它加速的是「怎么写」——查文献、组织语言、排版引用。核心判断（实验设计、数据解读）必须你自己把关。"),
    ("Q: 引用真的不会编造吗？",
     "A: 三层防护：Writer 只能引用文献库真实文献 → Verifier 拿到原文逐条比对 → Refiner 禁止为通过审查而删引用。疑似虚构会标记提醒。"),
    ("Q: 多方向文献混在一起影响准确吗？",
     "A: 检索时自动注入你的研究方向关键词 + 按分类精准过滤，热化学方向优先匹配热化学文献，不会被茶叶加工文献干扰。"),
    ("Q: 文献越来越多性能会下降吗？",
     "A: 当前 167 篇无压力。中期引入 HNSW 向量索引可支撑 10 万级，远期迁移专业向量数据库可支撑百万级。"),
]

for i, (q, a) in enumerate(qas):
    y = 1.5 + i * 1.15
    add_card(slide, 0.5, y, 12.3, 1.0, CARD_BG)
    add_text_box(slide, 0.7, y + 0.1, 11.8, 0.35, q, 16, GOLD, True)
    add_text_box(slide, 0.7, y + 0.45, 11.8, 0.45, a, 13, LIGHT)


# ═══════════════════════════════════════
# Slide 15: 致谢
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_divider(slide, 4.0, 2.2, 5.3, GOLD)
add_text_box(slide, 2.0, 1.0, 9.3, 1.0, "谢谢！", 64, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 2.0, 2.6, 9.3, 1.5,
    "论文智能助手 — 基于多 Agent 协作与 RAG 检索增强的学术写作平台\n欢迎提问和交流",
    18, LIGHT, False, PP_ALIGN.CENTER)
add_divider(slide, 4.0, 4.3, 5.3, GOLD)

add_text_box(slide, 2.0, 5.0, 9.3, 1.5,
    "技术栈: Next.js 16 + React + TypeScript + DeepSeek + 智谱 AI\n本地 RAG (BM25 + 向量语义) + Prisma + Turbopack",
    14, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)

# ── 保存 ──
output = "D:/project/论文助手/论文智能助手-项目演示.pptx"
prs.save(output)
print(f"✅ PPT 已生成: {output}")
print(f"   共 {len(prs.slides)} 页")
